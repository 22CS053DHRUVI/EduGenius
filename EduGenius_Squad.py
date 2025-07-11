import os
import fitz  # PyMuPDF
import requests
from bs4 import BeautifulSoup
from fpdf import FPDF
from pptx.api import Presentation
from pptx.util import Pt
from agno.agent import Agent
from agno.models.google import Gemini
import unicodedata

# Set your API key
from dotenv import load_dotenv
load_dotenv()

import os
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if GOOGLE_API_KEY is not None:
    os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY  # Replace with actual API key

# -------------- CONTENT EXTRACTORS --------------

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text

def fetch_text_from_url(url):
    headers = {"User-Agent": "Mozilla/5.0"}
    try:
        response = requests.get(url, headers=headers, timeout=10)
        soup = BeautifulSoup(response.text, "html.parser")
        paragraphs = soup.find_all('p')
        content = '\n'.join([para.get_text() for para in paragraphs])
        return content[:5000]
    except Exception as e:
        return f"Error fetching content from the URL: {e}"

# -------------- FILE SAVERS --------------

def clean_text_for_pdf(text):
    return text.encode("ascii", errors="ignore").decode()

def save_text_to_pdf(text, filename):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.set_font("Arial", size=12)

    text = clean_text_for_pdf(text)
    lines = text.split("\n")
    count = 1

    for line in lines:
        line = line.strip()
        if not line:
            pdf.ln()
            continue

        line = line.replace("**", "").replace("__", "").replace("*", "").replace("##", "")
        line = ''.join([c for c in line if ord(c) < 128])

        try:
            if line.endswith(":") or "Question" in line or "Q" in line[:3]:
                pdf.set_font("Arial", style='B', size=12)
                pdf.cell(0, 10, f"{count}. {line}", ln=True)
                count += 1
            else:
                pdf.set_font("Arial", style='', size=12)
                for segment in [line[i:i + 100] for i in range(0, len(line), 100)]:
                    pdf.multi_cell(0, 10, segment)
        except Exception as e:
            print(f"⚠️ Skipping problematic line: {line}\nError: {e}")
            continue

    try:
        pdf.output(filename)
        print(f"✅ PDF saved successfully: {filename}")
    except Exception as e:
        print(f"❌ Failed to save PDF: {e}")

def create_ppt_from_bullets(title, bullets, filename):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    max_bullets_per_slide = 5
    total_bullets = len(bullets)
    slide_count = 0

    for i in range(0, total_bullets, max_bullets_per_slide):
        current_bullets = bullets[i:i + max_bullets_per_slide]
        slide = prs.slides.add_slide(slide_layout)
        slide_count += 1
        slide_title = f"{title} - Slide {slide_count}" if slide_count > 1 else title
        slide.shapes.title.text = slide_title
        content = slide.placeholders[1]
        if hasattr(content, "text_frame"):
            tf = content.text_frame
            tf.clear()
        for point in current_bullets:
            if point.strip():
                p = tf.add_paragraph()
                p.text = point.strip()
                p.level = 0
                p.font.size = Pt(18)

    prs.save(filename)

# -------------- AGENTS --------------

doc_agent = Agent(name="DocuMaker", model=Gemini(id="gemini-2.0-flash-exp"), system_message=(
    "Generate only questions categorized as Easy, Medium, Hard. Each should be logical, scenario-based or topic-based. No instructions or extra formatting."
), markdown=True)

assign_agent = Agent(name="AssignPro", model=Gemini(id="gemini-2.0-flash"), system_message=(
    "You are an academic assignment designer.\n"
    "Generate assignment questions based on the provided topic or material.\n"
    "Include only well-framed questions that are:\n"
    "- Logical\n"
    "- Scenario-based\n"
    "- Topic-based\n"
    "Categorize them into three levels: Easy, Medium, and Hard.\n"
    "Do not include any instructions, formatting rules, or evaluation criteria—only questions."
), markdown=True)

ppt_agent = Agent(name="SlideGen", model=Gemini(id="gemini-2.0-flash"), instructions="""1. Organize content using clear bullet points under slide headings.
    2. Each slide should have a concise heading and 4–6 bullet points.
    3. Keep each bullet point short (max 20 words).
    4. Avoid complete paragraphs. Use phrases or short informative points.
    5. Group bullets by subtopic or concept for better visual segmentation.
    6. Use simple language and avoid jargon.

    Output format example:
    Slide 1: Introduction to Machine Learning
    - Definition of ML
    - History and evolution
...""")

paper_agent = Agent(name="PracticeEval", model=Gemini(id="gemini-2.0-flash"), system_message=(
    "You are an expert academic evaluator and paper generator. Based on the provided academic material, generate a well-structured practice test paper **based on user's chosen exam type**.\n\n"
    "🔹 If the user selects **30 marks**, assume it's for an internal exam:\n"
    "- Include ONLY **MCQs and short answer questions** (3-5 marks).\n"
    "- Mix logical, scenario-based, and conceptual questions.\n"
    "- Ensure total paper marks do not exceed 30.\n\n"
    "🔹 If the user selects **70 marks**, assume it's for an external/main exam:\n"
    "- Include a variety of questions: **MCQs, blanks, short answers (3/5 marks), long answers, thinking-based and scenario-based questions**.\n"
    "- Cover multiple difficulty levels (easy to hard).\n"
    "- Ensure the paper is properly structured and **totals to exactly 70 marks**.\n"
    "- Provide **answer keys and evaluation criteria**.\n\n"
    "Always:\n"
    "- Organize the paper section-wise if needed (Section A, B, etc.).\n"
    "- Label marks per question clearly.\n"
    "- Maintain academic tone and clean formatting."
), markdown=True)

# -------------- UTILITY PROMPT BUILDER --------------

def build_agent_prompt(content_type: str, content: str) -> str:
    prompt = (
        f"Please generate a high-quality {content_type} based on the academic content below.\n\n"
        "Follow these steps carefully:\n"
        "1. Understand the main topics and concepts in the text.\n"
        "2. Organize the content logically with clear structure.\n"
        "3. Include definitions, explanations, and relevant examples where appropriate.\n"
        "4. Use bullet points, headings, or numbered lists to improve readability.\n"
        "5. Keep the language clear, academic, and concise.\n"
        "6. Avoid repetition and irrelevant information.\n\n"
        "Here is the source content:\n\n"
        f"{content}\n\n"
        "Please start your output below:"
    )
    return prompt

# -------------- MAIN PROCESS FUNCTION --------------

def process_user_choice():
    print("What content would you like to generate?")
    print("1. Document\n2. Assignment\n3. PowerPoint\n4. Practice Paper")
    choice = input("Enter choice number: ").strip()

    print("\nChoose content source(s):")
    print("1. Upload PDF\n2. Enter URL\n(You can enter both separated by comma, e.g., 1,2)")
    sources = input("Enter 1, 2 or 1,2: ").strip().split(',')

    content = ""

    if '1' in sources:
        path = input("Enter path to PDF: ").strip()
        pdf_content = extract_text_from_pdf(path)
        content += pdf_content + "\n"

    if '2' in sources:
        url = input("Enter educational URL: ").strip()
        url_content = fetch_text_from_url(url)
        content += url_content + "\n"

    if not content.strip():
        print("❌ No valid content retrieved from selected sources.")
        return

    if choice == '1':
        prompt = build_agent_prompt("study notes", content)
        response = paper_agent.run(prompt)
        result = response.content.strip() if response and hasattr(response, "content") and response.content else ""
        if result:
            save_text_to_pdf(result, "output.pdf")
            print("✅ Paper saved to output.pdf")
        else:
            print("❌ Failed to generate paper.")

    elif choice == '2':
        prompt = build_agent_prompt("assignment", content)
        response = paper_agent.run(prompt)
        result = response.content.strip() if response and hasattr(response, "content") and response.content else ""
        if result:
            save_text_to_pdf(result, "Assignment.pdf")
            print("✅ Paper saved to Assignment.pdf")
        else:
            print("❌ Failed to generate paper.")

    elif choice == '3':
        prompt = build_agent_prompt("ppt content", content)
        result = ppt_agent.run(prompt).content.strip()
        bullets = [line.strip("-• ") for line in result.split("\n") if line.strip()]
        create_ppt_from_bullets("Generated Slides", bullets, "output.pptx")
        print("✅ Slides saved to output.pptx")

    elif choice == '4':
        print("Choose paper type:\n1. Internal Exam (30 marks)\n2. External Exam (70 marks)")
        ptype = input("Enter 1 or 2: ").strip()
        note = "Generate a 30 marks paper." if ptype == '1' else "Generate a 70 marks paper."
        prompt = f"{note}\n\n{build_agent_prompt('practice paper', content)}"
        response = paper_agent.run(prompt)
        result = response.content.strip() if response and hasattr(response, "content") and response.content else ""
        if result:
            save_text_to_pdf(result, "practice_paper.pdf")
            print("✅ Paper saved to practice_paper.pdf")
        else:
            print("❌ Failed to generate paper.")

    else:
        print("❌ Invalid choice!")

# -------------- MAIN EXECUTION --------------

if __name__ == "__main__":
    process_user_choice()
